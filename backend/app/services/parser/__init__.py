"""文档解析器：PDF / DOCX / PPTX → (章节列表, 每页文本)

设计要点（见 docs/01-architecture.md §4.1）：
- PDF 优先用 PyMuPDF（快，目录/文字/页码齐全）
- DOCX 用 python-docx（无 pandoc 依赖）
- PPTX 用 python-pptx
- 输出统一为 PageText 结构，方便后续切片与建索引
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class TocItem:
    """目录条目（从 PDF 书签或标题启发式提取）"""
    title: str
    level: int
    page: int  # 1-based


@dataclass
class ParseResult:
    """解析结果：章节树 + 每页文本"""
    pages: list[str] = field(default_factory=list)   # pages[i] = 第 i+1 页文本
    toc: list[TocItem] = field(default_factory=list)  # 有序目录
    total_pages: int = 0


class ParseError(Exception):
    pass


def parse_document(file_path: str | Path) -> ParseResult:
    """按扩展名分发解析。"""
    p = Path(file_path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(p)
    if suffix == ".docx":
        return _parse_docx(p)
    if suffix == ".pptx":
        return _parse_pptx(p)
    raise ParseError(f"不支持的文件类型: {suffix}")


# ---------- PDF ----------
def _parse_pdf(p: Path) -> ParseResult:
    try:
        import fitz  # PyMuPDF
    except ImportError as e:  # pragma: no cover
        raise ParseError("PyMuPDF 未安装，请运行 pip install pymupdf") from e

    doc = fitz.open(p)
    result = ParseResult(total_pages=doc.page_count)

    # 1. 书签目录
    raw_toc = doc.get_toc(simple=True)  # [(level, title, page), ...]
    for level, title, page in raw_toc:
        title = title.strip()
        if title:
            result.toc.append(TocItem(title=title, level=level, page=page))

    # 2. 逐页文本
    for page in doc:
        result.pages.append(page.get_text("text"))

    doc.close()
    return result


# ---------- DOCX ----------
# 标题样式映射：支持英文 Heading / 中文「大标题」「一级标题」等常见样式名
_HEADING_STYLE_MAP = {
    "heading 1": 1, "heading 2": 2, "heading 3": 3, "heading 4": 4, "heading 5": 5,
    "title": 1, "subtitle": 2,
    # 中文样式名（不同作者习惯）
    "大标题": 1, "标题": 1, "章标题": 1, "篇标题": 1,
    "一级标题": 1, "二级标题": 2, "三级标题": 3, "四级标题": 4, "五级标题": 5,
    "小节标题": 2,
}


def _docx_style_level(style_name: str) -> int | None:
    """把任意样式名解析为标题层级（None=非标题）。支持前缀匹配与数字后缀。"""
    if not style_name:
        return None
    name = style_name.strip()
    # 精确映射
    if name.lower() in _HEADING_STYLE_MAP:
        return _HEADING_STYLE_MAP[name.lower()]
    # 前缀匹配：Heading 1 / 标题 1 / heading1 等
    import re as _re
    m = _re.match(r"(?:heading|标题)\s*(\d)", name, _re.IGNORECASE)
    if m:
        lvl = int(m.group(1))
        return lvl if 1 <= lvl <= 6 else None
    # 中文前缀：标题 2 / 标题三 等
    m = _re.match(r"标题\s*([一二三四五六七八九十\d]+)", name)
    if m:
        _cn = {"一": 1, "二": 2, "三": 3, "四": 4, "五": 5, "六": 6}
        s = m.group(1)
        return _cn.get(s, int(s) if s.isdigit() else 1)
    return None


def _parse_docx(p: Path) -> ParseResult:
    try:
        import docx
    except ImportError as e:  # pragma: no cover
        raise ParseError("python-docx 未安装") from e

    d = docx.Document(p)
    result = ParseResult()

    page_text: list[str] = []
    current = []
    page_no = 1
    style_hit_count = 0  # 样式命中计数（判断是否走启发式回退）

    def flush():
        nonlocal current, page_no
        if current:
            # 段落间用空行分隔 (blank line separator)
            page_text.append("\n\n".join(current))
            current = []
            page_no += 1

    # 第一遍：样式标题识别
    for para in d.paragraphs:
        style = para.style.name if para.style is not None else ""
        text = para.text.strip()
        if not text:
            continue
        level = _docx_style_level(style)
        if level is not None:
            flush()
            result.toc.append(TocItem(title=text, level=level, page=page_no))
            style_hit_count += 1
        current.append(text)

    flush()

    # 第二遍：样式 + 启发式合并，并重建「每标题一页」的页结构
    from backend.app.services.rag.toc_heuristic import _CHAPTER_RE, _SECTION_RE
    import re as _re

    # 收集样式标题行号（避免启发式重复识别同一行）
    style_lines: set[int] = set()
    for idx, para in enumerate(d.paragraphs):
        style = para.style.name if para.style is not None else ""
        if _docx_style_level(style) is not None and para.text.strip():
            style_lines.add(idx)

    # 逐段扫描：识别标题（样式优先，其次启发式），重建 pages
    merged_toc: list[TocItem] = []
    seen_keys: set[str] = set()
    page_text = []
    current = []
    last_chapter_level = None  # 最近「第X章」是否出现（用于「一、」层级推断）

    def push_toc(title: str, level: int):
        nonlocal current
        key = title.strip()
        if not key or key in seen_keys:
            return
        seen_keys.add(key)
        if current:
            page_text.append("\n\n".join(current))
            current = []
        merged_toc.append(TocItem(title=_clean_docx_title(title), level=level, page=len(merged_toc) + 1))

    for idx, para in enumerate(d.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style is not None else ""
        level = _docx_style_level(style) if idx in style_lines else None
        if level is not None:
            push_toc(text, level)
            last_chapter_level = (level == 1)
            current.append(text)
            continue
        # 启发式：第X章 / 第X节
        m = _CHAPTER_RE.match(text)
        if m and len(text) <= 45:
            num = _re.search(r"第\s*([一二三四五六七八九十百千万0-9]+)", text)
            title = "第" + (num.group(1) if num else "?") + "章 " + _clean_docx_title(m.group(1))
            push_toc(title, 1)
            last_chapter_level = True
            current.append(text)
            continue
        m2 = _SECTION_RE.match(text)
        if m2 and len(text) <= 45:
            num = _re.search(r"第\s*([一二三四五六七八九十百千万0-9]+)", text)
            title = "第" + (num.group(1) if num else "?") + "节 " + _clean_docx_title(m2.group(1))
            push_toc(title, 2)
            current.append(text)
            continue
        # 中文序号「一、」：紧跟章后为二级，否则一级
        if _re.match(r"^[一二三四五六七八九十]+、", text) and 2 <= len(text) <= 40:
            push_toc(text, 2 if last_chapter_level else 1)
            current.append(text)
            continue
        # 其他情况：正文行
        current.append(text)

    if current:
        page_text.append("\n\n".join(current))
    result.toc = merged_toc
    result.pages = page_text or [""]
    result.total_pages = len(page_text)
    # 标题过少则放弃目录（保持整本一章）
    if len(result.toc) < 2:
        result.toc = []
    return result


def _clean_docx_title(raw: str) -> str:
    """清理标题尾部页码/年份/考试提示等噪声。"""
    import re as _re
    t = raw.strip()
    # 尾部纯数字/页码
    t = _re.sub(r"\s*[\d\s]+$", "", t)
    # 尾部「（4）」「(10)」等编号括号
    t = _re.sub(r"[（(]\d+[）)]\s*$", "", t)
    # 尾部「2012年简答：...」「2016论述：...」等考试提示
    t = _re.sub(r"20\d{2}\s*年?\s*(简答|论述|名词解释|计算|问答|考|预测).*$", "", t)
    t = _re.sub(r"（20\d{2}年?[^）]*）\s*$", "", t)
    # 中文间空格（字距）
    t = _re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fff])", "", t)
    return t.strip(" ，。、|｜:：")


# ---------- PPTX ----------
def _parse_pptx(p: Path) -> ParseResult:
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise ParseError("python-pptx 未安装") from e

    prs = Presentation(p)
    result = ParseResult(total_pages=len(prs.slides))
    page_texts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    texts.append(" | ".join(cells))
        page_texts.append("\n".join(texts))

    result.pages = page_texts
    return result